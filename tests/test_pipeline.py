#!/usr/bin/env python3
"""
InkDrop — End-to-End Test Pipeline
Validates all engines: processing, PDF insertion, DOCX insertion, PPTX insertion.
"""

import sys
import os
import tempfile
from pathlib import Path

# Add parent to path
sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))

PASS = "✓"
FAIL = "✗"
results = []


def test(name, fn):
    try:
        fn()
        results.append((name, True, None))
        print(f"  {PASS} {name}")
    except Exception as e:
        results.append((name, False, str(e)))
        print(f"  {FAIL} {name}: {e}")


def create_test_signature(path: str):
    """Create a minimal test signature image (dark strokes on white bg)."""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (400, 150), "white")
    draw = ImageDraw.Draw(img)
    # Draw a squiggle
    points = [(20, 75), (80, 30), (140, 110), (200, 50), (260, 90), (320, 40), (380, 75)]
    draw.line(points, fill="black", width=3)
    img.save(path)


def create_test_pdf(path: str):
    """Create a minimal test PDF with a signature line."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    c = canvas.Canvas(path, pagesize=letter)
    c.drawString(72, 700, "Test Document")
    c.drawString(72, 300, "Signature: ___________________________")
    c.save()


def create_test_docx(path: str):
    """Create a minimal test DOCX with a signature line."""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Test Document")
    doc.add_paragraph("This is the body of the document.")
    doc.add_paragraph("Signature: ___________________________")
    doc.save(path)


def create_test_pptx(path: str):
    """Create a minimal test PPTX with a signature placeholder."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agreement"
    slide.placeholders[1].text = "Signature: ___________________________"
    prs.save(path)


def run_tests():
    print("\n╔══════════════════════════════════════╗")
    print("║       InkDrop — Test Pipeline        ║")
    print("╚══════════════════════════════════════╝\n")

    with tempfile.TemporaryDirectory() as tmp:
        sig_raw = os.path.join(tmp, "sig_raw.png")
        sig_clean = os.path.join(tmp, "sig_clean.png")
        test_pdf = os.path.join(tmp, "test.pdf")
        test_docx = os.path.join(tmp, "test.docx")
        test_pptx = os.path.join(tmp, "test.pptx")
        out_pdf = os.path.join(tmp, "signed.pdf")
        out_docx = os.path.join(tmp, "signed.docx")
        out_pptx = os.path.join(tmp, "signed.pptx")

        # Create test fixtures
        create_test_signature(sig_raw)
        create_test_pdf(test_pdf)
        create_test_docx(test_docx)
        create_test_pptx(test_pptx)

        # Test 1: Signature Processing
        def t1():
            from process_signature import process_signature
            result = process_signature(sig_raw, sig_clean, threshold=200)
            assert Path(sig_clean).exists(), "Output not created"
            assert result["ink_coverage_percent"] > 0, "No ink detected"

        test("Process raw signature → clean PNG", t1)

        # Test 2: PDF Insertion
        def t2():
            from insert_pdf import insert_signature_pdf
            result = insert_signature_pdf(test_pdf, sig_clean, out_pdf, keyword="Signature")
            assert Path(out_pdf).exists(), "Output not created"
            assert result["page"] == 0

        test("Insert signature into PDF (keyword)", t2)

        # Test 3: DOCX Insertion
        def t3():
            from insert_docx import insert_signature_docx
            result = insert_signature_docx(test_docx, sig_clean, out_docx, keyword="Signature")
            assert Path(out_docx).exists(), "Output not created"

        test("Insert signature into DOCX (keyword)", t3)

        # Test 4: PPTX Insertion
        def t4():
            from insert_pptx import insert_signature_pptx
            result = insert_signature_pptx(test_pptx, sig_clean, out_pptx, keyword="Signature")
            assert Path(out_pptx).exists(), "Output not created"

        test("Insert signature into PPTX (keyword)", t4)

        # Test 5: Ink recoloring
        def t5():
            recolored = os.path.join(tmp, "sig_blue.png")
            from process_signature import process_signature
            result = process_signature(sig_raw, recolored, ink_color="#1a3c6e")
            assert Path(recolored).exists(), "Recolored output not created"

        test("Recolor ink to custom color", t5)

    # Summary
    passed = sum(1 for _, ok, _ in results if ok)
    total = len(results)
    print(f"\n{'═' * 40}")
    print(f"  Results: {passed}/{total} passed")
    if passed == total:
        print(f"  Status:  ALL CLEAR ✓")
    else:
        print(f"  Status:  {total - passed} FAILED")
        for name, ok, err in results:
            if not ok:
                print(f"    → {name}: {err}")
    print(f"{'═' * 40}\n")

    return 0 if passed == total else 1


if __name__ == "__main__":
    sys.exit(run_tests())
