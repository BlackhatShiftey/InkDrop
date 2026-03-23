#!/usr/bin/env python3
"""
InkDrop — PPTX Signature Insertion Engine
Drops a processed signature PNG into a PowerPoint presentation.

Supports:
  - Keyword-based placement (finds text on slides)
  - Slide number + position placement
  - Last-slide default placement

Usage:
    python insert_pptx.py deck.pptx signature.png output.pptx --find "Signature:"
    python insert_pptx.py deck.pptx signature.png output.pptx --slide 3 --x 2 --y 6
"""

import argparse
import sys
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("ERROR: python-pptx is required.")
    print("  pip install python-pptx --break-system-packages")
    sys.exit(1)


def find_signature_slide(prs: Presentation, keyword: str) -> dict | None:
    """Search slides for a keyword, return slide index and shape position."""
    patterns = [re.escape(keyword), r"_{5,}", r"\.{5,}", r"-{5,}"]

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                for pattern in patterns:
                    if re.search(pattern, text, re.IGNORECASE):
                        return {
                            "slide": slide_idx,
                            "x": shape.left,
                            "y": shape.top + shape.height,
                            "match": text[:60],
                        }
    return None


def insert_signature_pptx(
    pptx_path: str,
    sig_path: str,
    output_path: str,
    keyword: str = None,
    slide_num: int = None,
    x_inches: float = None,
    y_inches: float = None,
    sig_width_inches: float = 2.0,
) -> dict:
    """
    Insert a signature image into a PPTX presentation.

    Placement modes:
        1. keyword    — search for text marker on any slide
        2. slide_num  — specific slide (1-indexed) with optional x/y
        3. default    — last slide, lower-left area

    Returns:
        dict with placement metadata
    """
    prs = Presentation(pptx_path)

    target_slide_idx = (slide_num - 1) if slide_num else len(prs.slides) - 1
    target_x = Inches(x_inches) if x_inches else Inches(1.0)
    target_y = Inches(y_inches) if y_inches else Inches(6.0)

    if keyword:
        result = find_signature_slide(prs, keyword)
        if result:
            target_slide_idx = result["slide"]
            target_x = result["x"]
            target_y = result["y"]
        else:
            raise ValueError(f"Keyword '{keyword}' not found in presentation")

    slide = prs.slides[target_slide_idx]
    slide.shapes.add_picture(
        sig_path,
        target_x,
        target_y,
        width=Inches(sig_width_inches),
    )

    prs.save(output_path)

    return {
        "slide": target_slide_idx + 1,
        "position": f"({target_x}, {target_y})",
        "output": output_path,
    }


def main():
    parser = argparse.ArgumentParser(description="InkDrop — Insert signature into PPTX")
    parser.add_argument("pptx", help="Source PPTX file")
    parser.add_argument("signature", help="Processed signature PNG")
    parser.add_argument("output", help="Output PPTX path")
    parser.add_argument("--find", type=str, help="Keyword to locate signature area")
    parser.add_argument("--slide", type=int, help="Slide number (1-indexed)")
    parser.add_argument("--x", type=float, help="X position in inches")
    parser.add_argument("--y", type=float, help="Y position in inches")
    parser.add_argument("--width", type=float, default=2.0, help="Signature width in inches")

    args = parser.parse_args()

    for f in [args.pptx, args.signature]:
        if not Path(f).exists():
            print(f"ERROR: File not found: {f}")
            sys.exit(1)

    result = insert_signature_pptx(
        args.pptx, args.signature, args.output,
        keyword=args.find, slide_num=args.slide,
        x_inches=args.x, y_inches=args.y,
        sig_width_inches=args.width,
    )

    print(f"✓ InkDrop signed: slide {result['slide']}")
    print(f"  Output: {result['output']}")


if __name__ == "__main__":
    main()
